from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(248, 243, 239)
PANEL = RGBColor(255, 253, 249)
MUTE = RGBColor(111, 98, 108)
TITLE = RGBColor(59, 48, 56)
GOLD = RGBColor(200, 161, 90)
LAVENDER = RGBColor(184, 160, 216)
BLUE = RGBColor(175, 200, 222)
PINK = RGBColor(217, 168, 184)


def add_bg(slide, accent=GOLD):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    # soft decorative line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.color.rgb = accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.35), Inches(12.6), Inches(6.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(220, 210, 202)
    shape.shadow.inherit = False


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(6.5), Inches(1.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Songti SC'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = 'Heiti SC'
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTE


def add_card(slide, left, top, w, h, title=None, body=None, color=RGBColor(255, 255, 255)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(220, 210, 202)
    shape.shadow.inherit = False
    if title:
        tx = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), w - Inches(0.36), Inches(0.35))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Heiti SC'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TITLE
    if body:
        tx = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.45), w - Inches(0.36), h - Inches(0.55))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = 'Heiti SC'
            p.font.size = Pt(12)
            p.font.color.rgb = MUTE
            p.space_after = Pt(6)

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GOLD)
add_title(slide, '浅色梦幻图书馆设定集', '乙女游戏世界观 PPT / 视觉重构方向')
# big title
textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(7.2), Inches(1.1))
text_frame = textbox.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
p.text = 'Word Match × 乙女游戏世界观'
p.font.name = 'Songti SC'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = TITLE
p.alignment = PP_ALIGN.LEFT
p2 = text_frame.add_paragraph()
p2.text = '保留原有内容与原图，改造承载方式，打造设定集氛围。'
p2.font.name = 'Heiti SC'
p2.font.size = Pt(14)
p2.font.color.rgb = MUTE
# accent cards
add_card(slide, Inches(1.0), Inches(3.0), Inches(3.0), Inches(1.4), '核心关键词', ['月光', '纸页', '香槟金', '雾粉', '记忆碎片'], RGBColor(255, 255, 255))
add_card(slide, Inches(4.4), Inches(3.0), Inches(3.0), Inches(1.4), '视觉原则', ['浅色背景 / 层次感', '图片做插画卡承载', '文字进入奶白信息容器'], RGBColor(255, 255, 255))
add_card(slide, Inches(7.8), Inches(3.0), Inches(3.6), Inches(1.8), '目标效果', ['像乙女游戏设定集', '不是商务汇报', '有书页、信笺、命运线语言'], RGBColor(255, 255, 255))

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LAVENDER)
add_title(slide, '设计目标', '不改文字，不换图片内容；只改变承载方式与气质')
add_card(slide, Inches(0.9), Inches(1.6), Inches(3.7), Inches(2.1), '保留内容', ['原文字保留', '原图片保留', '只做视觉承载升级'], RGBColor(255,255,255))
add_card(slide, Inches(4.8), Inches(1.6), Inches(3.6), Inches(2.1), '改造重点', ['浅色纸感背景', '圆角相框与金边', '书签、信笺、记忆碎片符号'], RGBColor(255,255,255))
add_card(slide, Inches(8.7), Inches(1.6), Inches(3.2), Inches(2.1), '避坑', ['不要全页白遮罩', '不要把图洗灰', '不要用商务卡片硬替代氛围'], RGBColor(255,255,255))

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_title(slide, '色彩系统', '香槟金 / 雾粉 / 月光蓝 / 柔紫')
add_card(slide, Inches(0.9), Inches(1.7), Inches(2.8), Inches(2.4), '主背景', ['#F8F3EF', '象牙纸色'], RGBColor(255,255,255))
add_card(slide, Inches(4.0), Inches(1.7), Inches(2.8), Inches(2.4), '辅助背景', ['#F2EEF8', '淡紫雾色'], RGBColor(255,255,255))
add_card(slide, Inches(7.1), Inches(1.7), Inches(2.8), Inches(2.4), '信息卡底', ['#FFFDF9', '奶白卡片'], RGBColor(255,255,255))
add_card(slide, Inches(10.2), Inches(1.7), Inches(2.0), Inches(2.4), '强调色', ['#C8A15A', '#D9A8B8', '#AFC8DE'], RGBColor(255,255,255))

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PINK)
add_title(slide, '版式原则', '每页只做三件事：保留文字、保留图片、改变承载方式')
add_card(slide, Inches(0.9), Inches(1.8), Inches(3.8), Inches(2.6), '标题承载', ['左上书签条 + 香槟金细线', '标题偏宋体衬线，文学感更强'], RGBColor(255,255,255))
add_card(slide, Inches(4.9), Inches(1.8), Inches(3.8), Inches(2.6), '正文承载', ['半透明奶白卡片', '避免直接压在复杂图片上'], RGBColor(255,255,255))
add_card(slide, Inches(8.9), Inches(1.8), Inches(3.1), Inches(2.6), '图片承载', ['圆角 8–12px', '细金边 / 轻阴影', '像设定集插画卡'], RGBColor(255,255,255))

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GOLD)
add_title(slide, '核心机制页', '命运环、书页卡、记忆碎片与情感徽章')
add_card(slide, Inches(0.9), Inches(1.8), Inches(3.4), Inches(2.7), '命运环', ['化凡 / 邂逅 / 修复 / 选择', '四枚书签式节点'], RGBColor(255,255,255))
add_card(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(2.7), '情感机制', ['粘合 / 沟通 / 共鸣', '纵向递进的情感徽章'], RGBColor(255,255,255))
add_card(slide, Inches(8.7), Inches(1.8), Inches(3.2), Inches(2.7), '世界观语言', ['门、心、月、书签、丝带', '更像乙女游戏设定集'], RGBColor(255,255,255))

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LAVENDER)
add_title(slide, '结论', '将原始内容转成“设定集世界观”而不是“汇报页”')
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.2), Inches(3.1))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255,255,255); shape.line.color.rgb = GOLD
textbox = slide.shapes.add_textbox(Inches(1.35), Inches(2.1), Inches(10.6), Inches(2.35))
text_frame = textbox.text_frame
text_frame.word_wrap = True
for i, line in enumerate([
    '这套方案的核心不是把颜色变得更“漂亮”，而是把每一页都变成“乙女游戏设定集”的承载形式。',
    '图片保留原内容，文字保留原信息，页面背景、卡片、符号与线条共同构成世界观感。',
    '这样既能保留原始素材，又能把它转成更有氛围的演示文稿。'
]):
    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
    p.text = line
    p.font.name = 'Heiti SC'
    p.font.size = Pt(18 if i == 0 else 16)
    p.font.color.rgb = TITLE
    p.space_after = Pt(6)

prs.save('乙女游戏世界观设定集_PPT.pptx')
print('已生成: 乙女游戏世界观设定集_PPT.pptx')
